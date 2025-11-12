"""文档解析器模块

根据不同文件类型提供不同的解析方法，统一输出 Markdown 格式文本
支持的文件类型: txt, md, pdf, docx, doc, html, htm, xlsx, xls, csv, pptx, jpg, png
"""

import asyncio
from io import BytesIO
from pathlib import Path
from typing import Any

from loguru import logger


class BaseParser:
    """解析器基类"""

    @staticmethod
    async def parse(file_content: bytes, filename: str, **kwargs) -> str:
        """解析文档为 Markdown 格式文本

        Args:
            file_content: 文件内容（字节）
            filename: 文件名
            **kwargs: 额外参数

        Returns:
            str: Markdown 格式的文档内容
        """
        raise NotImplementedError("子类必须实现 parse 方法")


class TextParser(BaseParser):
    """纯文本解析器 (txt, md)"""

    @staticmethod
    async def parse(file_content: bytes, filename: str, **kwargs) -> str:
        """解析纯文本文件"""
        try:
            # 尝试多种编码
            encodings = ["utf-8", "gbk", "gb2312", "gb18030", "latin1"]

            for encoding in encodings:
                try:
                    text = file_content.decode(encoding)
                    logger.info(f"使用 {encoding} 编码成功解析文件: {filename}")
                    return text
                except UnicodeDecodeError:
                    continue

            # 如果所有编码都失败，使用 utf-8 并忽略错误
            text = file_content.decode("utf-8", errors="ignore")
            logger.warning(f"使用 utf-8 忽略错误解析文件: {filename}")
            return text

        except Exception as e:
            logger.error(f"文本解析失败 {filename}: {e}")
            raise ValueError(f"文本解析失败: {e}")


class PDFParser(BaseParser):
    """PDF 解析器"""

    @staticmethod
    async def parse(
        file_content: bytes, filename: str, use_ocr: bool = False, ocr_config: dict[str, Any] | None = None, **kwargs
    ) -> str:
        """解析 PDF 文档

        Args:
            file_content: 文件内容
            filename: 文件名
            use_ocr: 是否使用 OCR（适用于图片型 PDF）
            ocr_config: OCR 配置
                - model: OCR 模型名称，默认 'deepseek-ai/DeepSeek-OCR'
                - dpi: PDF 转图片的 DPI，默认 144
                - max_concurrent: 最大并发数，默认 30
                - extract_images: 是否提取图片，默认 False

        Returns:
            str: Markdown 格式的文档内容
        """
        try:
            if use_ocr:
                # 使用 OCR 解析（适用于扫描版 PDF）
                from ragserver.app.utils.deepseek_ocr import DocumentOCRParser

                ocr_config = ocr_config or {}
                parser = DocumentOCRParser(
                    model=ocr_config.get("model", "deepseek-ai/DeepSeek-OCR"),
                    dpi=ocr_config.get("dpi", 144),
                    max_concurrent=ocr_config.get("max_concurrent", 30),
                )

                # 保存临时文件
                import tempfile

                with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp_file:
                    tmp_file.write(file_content)
                    tmp_path = tmp_file.name

                try:
                    text = await parser.parse_document(tmp_path, extract_images=ocr_config.get("extract_images", False))
                    logger.info(f"使用 OCR 解析 PDF 成功: {filename}")
                    return text
                finally:
                    # 清理临时文件
                    import os

                    if os.path.exists(tmp_path):
                        os.unlink(tmp_path)

            else:
                # 使用标准文本提取（适用于文本型 PDF）
                from pdfminer.high_level import extract_text

                file_obj = BytesIO(file_content)
                text = await asyncio.to_thread(extract_text, file_obj)

                # 检查是否提取到有效内容
                if not text or len(text.strip()) < 50:
                    logger.warning(f"PDF 文本提取内容过少，建议使用 OCR: {filename}")
                    # 可以自动切换到 OCR
                    # return await PDFParser.parse(file_content, filename, use_ocr=True, ocr_config=ocr_config)

                logger.info(f"PDF 解析成功: {filename}, 长度: {len(text)}")
                return text

        except Exception as e:
            logger.error(f"PDF 解析失败 {filename}: {e}")
            raise ValueError(f"PDF 解析失败: {e}")


class DocxParser(BaseParser):
    """DOCX 文档解析器"""

    @staticmethod
    async def parse(file_content: bytes, filename: str, **kwargs) -> str:
        """解析 DOCX 文档"""
        try:
            import docx

            file_obj = BytesIO(file_content)
            doc = await asyncio.to_thread(docx.Document, file_obj)

            # 提取段落
            paragraphs = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    # 根据段落样式添加 Markdown 标记
                    style = para.style.name.lower()
                    if "heading 1" in style:
                        paragraphs.append(f"# {text}")
                    elif "heading 2" in style:
                        paragraphs.append(f"## {text}")
                    elif "heading 3" in style:
                        paragraphs.append(f"### {text}")
                    else:
                        paragraphs.append(text)

            # 提取表格
            tables = []
            for table in doc.tables:
                table_md = []
                for i, row in enumerate(table.rows):
                    cells = [cell.text.strip() for cell in row.cells]
                    table_md.append("| " + " | ".join(cells) + " |")
                    # 添加表头分隔符
                    if i == 0:
                        table_md.append("| " + " | ".join(["---"] * len(cells)) + " |")

                if table_md:
                    tables.append("\n".join(table_md))

            # 合并段落和表格
            content = "\n\n".join(paragraphs)
            if tables:
                content += "\n\n" + "\n\n".join(tables)

            logger.info(f"DOCX 解析成功: {filename}, 段落数: {len(paragraphs)}, 表格数: {len(tables)}")
            return content

        except Exception as e:
            logger.error(f"DOCX 解析失败 {filename}: {e}")
            raise ValueError(f"DOCX 解析失败: {e}")


class HTMLParser(BaseParser):
    """HTML 解析器"""

    @staticmethod
    async def parse(file_content: bytes, filename: str, preserve_links: bool = True, **kwargs) -> str:
        """解析 HTML 文档

        Args:
            file_content: 文件内容
            filename: 文件名
            preserve_links: 是否保留链接（转为 Markdown 格式）

        Returns:
            str: Markdown 格式的文档内容
        """
        try:
            from bs4 import BeautifulSoup

            soup = BeautifulSoup(file_content, "html.parser")

            # 移除 script、style、nav、footer 等标签
            for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
                tag.decompose()

            # 转换标题标签
            for i in range(1, 7):
                for heading in soup.find_all(f"h{i}"):
                    heading.string = f"{'#' * i} {heading.get_text()}"

            # 转换链接
            if preserve_links:
                for link in soup.find_all("a"):
                    href = link.get("href", "")
                    text = link.get_text()
                    if href and text:
                        link.string = f"[{text}]({href})"

            # 转换列表
            for ul in soup.find_all("ul"):
                for li in ul.find_all("li"):
                    li.string = f"- {li.get_text()}"

            for ol in soup.find_all("ol"):
                for idx, li in enumerate(ol.find_all("li"), 1):
                    li.string = f"{idx}. {li.get_text()}"

            # 获取文本
            text = soup.get_text()

            # 清理多余空行
            lines = [line.strip() for line in text.splitlines() if line.strip()]
            content = "\n\n".join(lines)

            logger.info(f"HTML 解析成功: {filename}, 长度: {len(content)}")
            return content

        except Exception as e:
            logger.error(f"HTML 解析失败 {filename}: {e}")
            raise ValueError(f"HTML 解析失败: {e}")


class ExcelParser(BaseParser):
    """Excel/CSV 解析器"""

    @staticmethod
    async def parse(file_content: bytes, filename: str, file_type: str = "xlsx", **kwargs) -> str:
        """解析 Excel/CSV 文档

        Args:
            file_content: 文件内容
            filename: 文件名
            file_type: 文件类型 (xlsx, xls, csv)

        Returns:
            str: Markdown 格式的文档内容
        """
        try:
            import pandas as pd

            file_obj = BytesIO(file_content)

            if file_type == "csv":
                # 读取 CSV
                df = await asyncio.to_thread(pd.read_csv, file_obj)
                markdown = df.to_markdown(index=False)
                logger.info(f"CSV 解析成功: {filename}, 行数: {len(df)}")
                return markdown

            else:
                # 读取 Excel 的所有 sheet
                sheets = await asyncio.to_thread(
                    pd.read_excel, file_obj, sheet_name=None, engine="openpyxl" if file_type == "xlsx" else "xlrd"
                )

                if not isinstance(sheets, dict):
                    sheets = {"Sheet1": sheets}

                # 将所有 sheet 转换为 Markdown
                text_parts = []
                for sheet_name, df in sheets.items():
                    text_parts.append(f"# Sheet: {sheet_name}\n")

                    # 跳过空 sheet
                    if df.empty:
                        text_parts.append("_（此工作表为空）_\n")
                        continue

                    # 转换为 Markdown 表格
                    markdown = df.to_markdown(index=False)
                    text_parts.append(markdown)
                    text_parts.append("\n")

                content = "\n".join(text_parts)
                logger.info(f"Excel 解析成功: {filename}, Sheet数: {len(sheets)}")
                return content

        except Exception as e:
            logger.error(f"Excel 解析失败 {filename}: {e}")
            raise ValueError(f"Excel 解析失败: {e}")


class PPTXParser(BaseParser):
    """PPTX 演示文稿解析器"""

    @staticmethod
    async def parse(file_content: bytes, filename: str, **kwargs) -> str:
        """解析 PPTX 文档"""
        try:
            from pptx import Presentation

            file_obj = BytesIO(file_content)
            prs = await asyncio.to_thread(Presentation, file_obj)

            slides_content = []

            for idx, slide in enumerate(prs.slides, 1):
                slide_parts = [f"# Slide {idx}\n"]

                # 提取幻灯片中的所有文本
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        # 标题通常是第一个文本框
                        if shape.is_placeholder and shape.placeholder_format.type == 1:
                            slide_parts.append(f"## {text}\n")
                        else:
                            slide_parts.append(f"{text}\n")

                    # 提取表格
                    if shape.has_table:
                        table = shape.table
                        table_md = []
                        for row_idx, row in enumerate(table.rows):
                            cells = [cell.text.strip() for cell in row.cells]
                            table_md.append("| " + " | ".join(cells) + " |")
                            if row_idx == 0:
                                table_md.append("| " + " | ".join(["---"] * len(cells)) + " |")

                        slide_parts.append("\n".join(table_md) + "\n")

                slides_content.append("\n".join(slide_parts))

            content = "\n\n---\n\n".join(slides_content)
            logger.info(f"PPTX 解析成功: {filename}, 幻灯片数: {len(prs.slides)}")
            return content

        except Exception as e:
            logger.error(f"PPTX 解析失败 {filename}: {e}")
            raise ValueError(f"PPTX 解析失败: {e}")


class ImageParser(BaseParser):
    """图片解析器 (jpg, png)"""

    @staticmethod
    async def parse(file_content: bytes, filename: str, ocr_config: dict[str, Any] | None = None, **kwargs) -> str:
        """解析图片文档（使用 OCR）

        Args:
            file_content: 文件内容
            filename: 文件名
            ocr_config: OCR 配置
                - model: OCR 模型名称，默认 'deepseek-ai/DeepSeek-OCR'
                - max_concurrent: 最大并发数，默认 30

        Returns:
            str: Markdown 格式的文档内容
        """
        try:
            from ragserver.app.utils.deepseek_ocr import DocumentOCRParser

            ocr_config = ocr_config or {}
            parser = DocumentOCRParser(
                model=ocr_config.get("model", "deepseek-ai/DeepSeek-OCR"),
                max_concurrent=ocr_config.get("max_concurrent", 30),
            )

            # 保存临时文件
            import tempfile

            suffix = Path(filename).suffix
            with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp_file:
                tmp_file.write(file_content)
                tmp_path = tmp_file.name

            try:
                text = await parser.parse_document(tmp_path)
                logger.info(f"图片 OCR 解析成功: {filename}")
                return text
            finally:
                # 清理临时文件
                import os

                if os.path.exists(tmp_path):
                    os.unlink(tmp_path)

        except Exception as e:
            logger.error(f"图片解析失败 {filename}: {e}")
            raise ValueError(f"图片解析失败: {e}")


# ==================== 解析器工厂 ====================


class DocumentParserFactory:
    """文档解析器工厂"""

    # 文件类型到解析器的映射
    PARSER_MAP = {
        "txt": TextParser,
        "md": TextParser,
        "pdf": PDFParser,
        "docx": DocxParser,
        "doc": DocxParser,
        "html": HTMLParser,
        "htm": HTMLParser,
        "xlsx": ExcelParser,
        "xls": ExcelParser,
        "csv": ExcelParser,
        "pptx": PPTXParser,
        "jpg": ImageParser,
        "jpeg": ImageParser,
        "png": ImageParser,
    }

    @classmethod
    def get_parser(cls, file_type: str) -> BaseParser:
        """根据文件类型获取对应的解析器

        Args:
            file_type: 文件类型（不含点号）

        Returns:
            BaseParser: 解析器实例

        Raises:
            ValueError: 不支持的文件类型
        """
        file_type = file_type.lower().strip(".")

        parser_class = cls.PARSER_MAP.get(file_type)
        if not parser_class:
            raise ValueError(f"不支持的文件类型: {file_type}。支持的类型: {', '.join(cls.PARSER_MAP.keys())}")

        return parser_class()

    @classmethod
    async def parse_document(cls, file_content: bytes, file_type: str, filename: str, **kwargs) -> str:
        """统一的文档解析接口

        Args:
            file_content: 文件内容（字节）
            file_type: 文件类型（pdf, docx, txt等）
            filename: 文件名
            **kwargs: 额外参数
                - use_ocr: 是否使用 OCR（PDF）
                - ocr_config: OCR 配置
                - preserve_links: 是否保留链接（HTML）

        Returns:
            str: Markdown 格式的文档内容

        Raises:
            ValueError: 解析失败或不支持的文件类型
        """
        try:
            parser = cls.get_parser(file_type)

            # 传递文件类型（Excel需要）
            if file_type in ["xlsx", "xls", "csv"]:
                kwargs["file_type"] = file_type

            text = await parser.parse(file_content, filename, **kwargs)

            if not text or not text.strip():
                raise ValueError(f"解析结果为空: {filename}")

            return text

        except Exception as e:
            logger.error(f"文档解析失败 {filename} ({file_type}): {e}")
            raise


# ==================== 便捷函数 ====================


async def parse_document(
    file_content: bytes,
    file_type: str,
    filename: str,
    use_ocr: bool = False,
    ocr_config: dict[str, Any] | None = None,
    **kwargs,
) -> str:
    """解析文档的便捷函数

    Args:
        file_content: 文件内容（字节）
        file_type: 文件类型（pdf, docx, txt等）
        filename: 文件名
        use_ocr: 是否使用 OCR（适用于 PDF 和图片）
        ocr_config: OCR 配置
        **kwargs: 其他参数

    Returns:
        str: Markdown 格式的文档内容
    """
    if use_ocr:
        kwargs["use_ocr"] = True
        kwargs["ocr_config"] = ocr_config

    return await DocumentParserFactory.parse_document(
        file_content=file_content, file_type=file_type, filename=filename, **kwargs
    )


def get_supported_file_types() -> list:
    """获取支持的文件类型列表"""
    return list(DocumentParserFactory.PARSER_MAP.keys())


def is_supported_file_type(file_type: str) -> bool:
    """检查文件类型是否支持"""
    file_type = file_type.lower().strip(".")
    return file_type in DocumentParserFactory.PARSER_MAP


# ==================== 示例用法 ====================

if __name__ == "__main__":
    # 示例：解析不同类型的文档
    async def test_parsers():
        # 1. 解析 PDF（文本型）
        with open("example.pdf", "rb") as f:
            content = await parse_document(f.read(), "pdf", "example.pdf")
            print(f"PDF 内容长度: {len(content)}")

        # 2. 解析 PDF（扫描型，使用 OCR）
        with open("scanned.pdf", "rb") as f:
            content = await parse_document(f.read(), "pdf", "scanned.pdf", use_ocr=True, ocr_config={"dpi": 144})
            print(f"OCR PDF 内容长度: {len(content)}")

        # 3. 解析 DOCX
        with open("example.docx", "rb") as f:
            content = await parse_document(f.read(), "docx", "example.docx")
            print(f"DOCX 内容长度: {len(content)}")

        # 4. 解析图片（OCR）
        with open("image.jpg", "rb") as f:
            content = await parse_document(f.read(), "jpg", "image.jpg")
            print(f"图片 OCR 内容长度: {len(content)}")

    # 运行测试
    # asyncio.run(test_parsers())

    # 查看支持的文件类型
    print(f"支持的文件类型: {get_supported_file_types()}")
